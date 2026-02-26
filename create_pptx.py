#!/usr/bin/env python3
"""
BioRadio Hackathon 2026 - Kickoff Presentation Generator
Generates a professional 16:9 PowerPoint presentation using python-pptx.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

# ── Color Palette ──────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(20, 40, 80)
MEDIUM_BLUE = RGBColor(30, 60, 110)
TEAL = RGBColor(0, 180, 200)
LIGHT_TEAL = RGBColor(0, 210, 230)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(240, 244, 248)
LIGHT_GRAY = RGBColor(220, 225, 230)
DARK_TEXT = RGBColor(30, 30, 50)
MEDIUM_GRAY = RGBColor(100, 110, 130)
TABLE_HEADER_BG = RGBColor(20, 40, 80)
TABLE_ALT_ROW = RGBColor(235, 245, 250)
TABLE_BORDER = RGBColor(180, 200, 220)
ACCENT_ORANGE = RGBColor(255, 165, 0)

# ── Dimensions (16:9) ─────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_NAME = "Calibri"

# ── Presentation Setup ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout for full control
BLANK_LAYOUT = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════

def add_dark_background(slide):
    """Fill entire slide with dark blue background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE


def add_light_background(slide):
    """Fill entire slide with off-white background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = OFF_WHITE


def add_accent_bar(slide, top=0, height=Inches(0.08), color=TEAL):
    """Add a horizontal teal accent bar across the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, top, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bottom_bar(slide, color=TEAL, height=Inches(0.06)):
    """Add a thin accent bar at the bottom of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_HEIGHT - height, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_side_accent(slide, color=TEAL, width=Inches(0.12)):
    """Add a vertical accent bar on the left side."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, width, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_textbox(slide, text, left=Inches(0.8), top=Inches(0.4),
                      width=Inches(11.5), height=Inches(0.9),
                      font_size=36, font_color=DARK_BLUE, bold=True):
    """Add a title text box with consistent formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_color
    p.font.name = FONT_NAME
    return txBox


def add_body_textbox(slide, left, top, width, height):
    """Add a text box and return the text_frame for further editing."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def set_paragraph(p, text, size=20, color=DARK_TEXT, bold=False,
                  italic=False, alignment=PP_ALIGN.LEFT, space_after=Pt(8),
                  space_before=Pt(0), font_name=FONT_NAME):
    """Configure a paragraph with formatting."""
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before


def add_bullet(tf, text, size=20, color=DARK_TEXT, bold=False, level=0,
               space_after=Pt(6), bullet_char="\u2022"):
    """Add a bullet point paragraph to a text frame."""
    p = tf.add_paragraph()
    p.text = f"{bullet_char} {text}" if level == 0 else f"   {bullet_char} {text}"
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT_NAME
    p.space_after = space_after
    p.space_before = Pt(2)
    p.level = level
    return p


def add_numbered_item(tf, number, title, description, title_size=20,
                      desc_size=18, title_color=DARK_BLUE,
                      desc_color=MEDIUM_GRAY):
    """Add a numbered item with title and description."""
    p = tf.add_paragraph()
    run_num = p.add_run()
    run_num.text = f"{number}. "
    run_num.font.size = Pt(title_size)
    run_num.font.color.rgb = TEAL
    run_num.font.bold = True
    run_num.font.name = FONT_NAME

    run_title = p.add_run()
    run_title.text = title
    run_title.font.size = Pt(title_size)
    run_title.font.color.rgb = title_color
    run_title.font.bold = True
    run_title.font.name = FONT_NAME

    if description:
        run_desc = p.add_run()
        run_desc.text = f" \u2014 {description}"
        run_desc.font.size = Pt(desc_size)
        run_desc.font.color.rgb = desc_color
        run_desc.font.bold = False
        run_desc.font.name = FONT_NAME

    p.space_after = Pt(10)
    p.space_before = Pt(4)
    return p


def format_table_cell(cell, text, font_size=16, bold=False, color=DARK_TEXT,
                      alignment=PP_ALIGN.LEFT, bg_color=None,
                      font_name=FONT_NAME):
    """Format a table cell with text and styling."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Set margins
    cell.margin_left = Inches(0.12)
    cell.margin_right = Inches(0.12)
    cell.margin_top = Inches(0.06)
    cell.margin_bottom = Inches(0.06)

    if bg_color:
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = bg_color


def set_cell_border(cell, color=TABLE_BORDER, width=Pt(0.75)):
    """Set borders on a table cell using low-level XML manipulation."""
    from pptx.oxml.ns import qn
    from lxml import etree

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
        # Remove existing border if present
        existing = tcPr.find(qn(f'a:{border_name}'))
        if existing is not None:
            tcPr.remove(existing)

        ln = etree.SubElement(tcPr, qn(f'a:{border_name}'))
        ln.set('w', str(int(width)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')

        solidFill = etree.SubElement(ln, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', f'{color.red:02X}{color.green:02X}{color.blue:02X}' if hasattr(color, 'red') else str(color))

        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'solid')


def create_styled_table(slide, rows, cols, left, top, width, height,
                        headers, data, col_widths=None):
    """Create a professionally styled table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Format header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        format_table_cell(cell, header, font_size=16, bold=True,
                         color=WHITE, bg_color=TABLE_HEADER_BG,
                         alignment=PP_ALIGN.CENTER)
        set_cell_border(cell, color=RGBColor(15, 30, 60))

    # Format data rows
    for i, row_data in enumerate(data):
        bg = TABLE_ALT_ROW if i % 2 == 0 else WHITE
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i + 1, j)
            format_table_cell(cell, cell_text, font_size=15, bold=False,
                            color=DARK_TEXT, bg_color=bg)
            set_cell_border(cell, color=TABLE_BORDER)

    return table_shape


def add_rounded_rect_card(slide, left, top, width, height, fill_color=WHITE,
                          border_color=TEAL, border_width=Pt(1.5)):
    """Add a rounded rectangle card shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_slide_number(slide, number, total=14):
    """Add a slide number indicator at the bottom right."""
    txBox = slide.shapes.add_textbox(
        Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number} / {total}"
    p.font.size = Pt(11)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.RIGHT


def content_slide_setup(slide, title_text, slide_num):
    """Common setup for content slides: light bg, side accent, title, bars."""
    add_light_background(slide)
    add_side_accent(slide, color=TEAL, width=Inches(0.1))
    add_accent_bar(slide, top=Inches(1.25), height=Inches(0.04), color=TEAL)
    add_bottom_bar(slide, color=DARK_BLUE, height=Inches(0.05))
    add_title_textbox(slide, title_text, left=Inches(0.6), top=Inches(0.35),
                      font_size=34, font_color=DARK_BLUE)
    add_slide_number(slide, slide_num)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════════════════
def create_slide_1():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_dark_background(slide)

    # Top decorative bar
    add_accent_bar(slide, top=0, height=Inches(0.1), color=TEAL)

    # Decorative teal rectangle on the left
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.0), Inches(0.35), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    # Subtle decorative element - right side
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(12.5), Inches(1.0), Inches(0.8), Inches(0.08)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = TEAL
    shape2.line.fill.background()

    shape3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(11.8), Inches(1.2), Inches(1.5), Inches(0.08)
    )
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(0, 140, 160)
    shape3.line.fill.background()

    # Main title
    txBox = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.2), Inches(11), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "BioRadio"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = TEAL
    run.font.name = FONT_NAME

    run2 = p.add_run()
    run2.text = " Hackathon 2026"
    run2.font.size = Pt(54)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    run2.font.name = FONT_NAME

    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.7), Inches(10), Inches(0.8)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "From Biosignals to Real-Time Control"
    p2.font.size = Pt(28)
    p2.font.color.rgb = LIGHT_TEAL
    p2.font.name = FONT_NAME
    p2.font.italic = True
    p2.alignment = PP_ALIGN.CENTER

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.7), Inches(4.3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Footer
    txBox3 = slide.shapes.add_textbox(
        Inches(2), Inches(5.1), Inches(9), Inches(0.6)
    )
    tf3 = txBox3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]

    run_badge = p3.add_run()
    run_badge.text = "\u23F1  "
    run_badge.font.size = Pt(22)
    run_badge.font.color.rgb = TEAL
    run_badge.font.name = FONT_NAME

    run_footer = p3.add_run()
    run_footer.text = "24-Hour Build Challenge"
    run_footer.font.size = Pt(24)
    run_footer.font.color.rgb = RGBColor(180, 200, 220)
    run_footer.font.bold = True
    run_footer.font.name = FONT_NAME
    p3.alignment = PP_ALIGN.CENTER

    # Bottom decorative bar
    add_accent_bar(slide, top=SLIDE_HEIGHT - Inches(0.1),
                   height=Inches(0.1), color=TEAL)

    add_slide_number(slide, 1)
    # Override slide number color for dark bg
    slide.shapes[-1].text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 120, 150)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2: What You'll Build
# ══════════════════════════════════════════════════════════════════════════
def create_slide_2():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "What You'll Build", 2)

    # Pipeline diagram - create 4 boxes with arrows
    box_labels = ["Biosignals", "Feature\nExtraction", "ML\nClassifier", "Control\nSystem"]
    box_subtitles = ["(BioRadio)", "(your code)", "(trained model)", "(your choice)"]
    box_width = Inches(2.4)
    box_height = Inches(1.6)
    start_x = Inches(0.7)
    box_y = Inches(2.0)
    gap = Inches(0.85)

    for i, (label, subtitle) in enumerate(zip(box_labels, box_subtitles)):
        x = start_x + i * (box_width + gap)

        # Box shape
        card = add_rounded_rect_card(
            slide, x, box_y, box_width, box_height,
            fill_color=WHITE if i > 0 else RGBColor(230, 248, 250),
            border_color=TEAL if i == 0 else RGBColor(180, 200, 220)
        )

        # Label
        txBox = slide.shapes.add_textbox(
            x + Inches(0.1), box_y + Inches(0.2),
            box_width - Inches(0.2), Inches(0.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        txBox2 = slide.shapes.add_textbox(
            x + Inches(0.1), box_y + Inches(1.0),
            box_width - Inches(0.2), Inches(0.5)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = MEDIUM_GRAY
        p2.font.name = FONT_NAME
        p2.font.italic = True
        p2.alignment = PP_ALIGN.CENTER

        # Arrow between boxes (not after the last one)
        if i < 3:
            arrow_x = x + box_width + Inches(0.1)
            arrow_y = box_y + box_height / 2 - Inches(0.15)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y,
                Inches(0.6), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEAL
            arrow.line.fill.background()

    # Description text below
    tf = add_body_textbox(
        slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.5)
    )
    p = tf.paragraphs[0]
    p.text = ""

    # Main emphasis line
    p2 = tf.add_paragraph()
    run1 = p2.add_run()
    run1.text = "Collect biosignal data"
    run1.font.size = Pt(22)
    run1.font.bold = True
    run1.font.color.rgb = TEAL
    run1.font.name = FONT_NAME

    run2 = p2.add_run()
    run2.text = ", "
    run2.font.size = Pt(22)
    run2.font.color.rgb = DARK_TEXT
    run2.font.name = FONT_NAME

    run3 = p2.add_run()
    run3.text = "train a classifier"
    run3.font.size = Pt(22)
    run3.font.bold = True
    run3.font.color.rgb = TEAL
    run3.font.name = FONT_NAME

    run4 = p2.add_run()
    run4.text = ", and use the output to "
    run4.font.size = Pt(22)
    run4.font.color.rgb = DARK_TEXT
    run4.font.name = FONT_NAME

    run5 = p2.add_run()
    run5.text = "control something in real-time"
    run5.font.size = Pt(22)
    run5.font.bold = True
    run5.font.color.rgb = TEAL
    run5.font.name = FONT_NAME

    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3: Signal Types Available
# ══════════════════════════════════════════════════════════════════════════
def create_slide_3():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Signal Types Available", 3)

    headers = ["Signal", "What It Measures", "Typical Use"]
    data = [
        ["EMG", "Muscle electrical activity", "Gesture control, prosthetics"],
        ["EOG", "Eye movement & blinks", "Gaze control, blink detection"],
        ["EEG", "Brain electrical activity", "BCI, attention monitoring"],
        ["GSR", "Skin conductance", "Stress detection, arousal"],
        ["IMU", "Acceleration & rotation", "Motion tracking"],
    ]

    col_widths = [Inches(1.8), Inches(4.0), Inches(5.0)]

    create_styled_table(
        slide, rows=6, cols=3,
        left=Inches(0.9), top=Inches(1.7),
        width=Inches(10.8), height=Inches(4.5),
        headers=headers, data=data, col_widths=col_widths
    )


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4: The Hackathon Pipeline (detailed)
# ══════════════════════════════════════════════════════════════════════════
def create_slide_4():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "The Hackathon Pipeline", 4)

    steps = [
        ("Data Collection",
         "Use the hackathon GUI to record labeled training data"),
        ("Feature Extraction",
         "Extract features (RMS, frequency bands, zero crossings, etc.)"),
        ("Model Training",
         "Train a classifier (scikit-learn, PyTorch, TensorFlow, etc.)"),
        ("Real-Time Inference",
         "Run your model on live BioRadio data"),
        ("Control",
         "Map classifier outputs to actions in your chosen system"),
    ]

    start_y = Inches(1.65)
    step_height = Inches(1.0)

    for i, (title, desc) in enumerate(steps):
        y = start_y + i * step_height

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.1),
            Inches(0.55), Inches(0.55)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()

        # Number text
        tf_circle = circle.text_frame
        tf_circle.word_wrap = False
        p_num = tf_circle.paragraphs[0]
        p_num.text = str(i + 1)
        p_num.font.size = Pt(22)
        p_num.font.bold = True
        p_num.font.color.rgb = WHITE
        p_num.font.name = FONT_NAME
        p_num.alignment = PP_ALIGN.CENTER
        tf_circle.vertical_anchor = MSO_ANCHOR.MIDDLE  # type: ignore

        # Connecting line between circles (not after last)
        if i < 4:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1.05), y + Inches(0.66),
                Inches(0.04), Inches(0.45)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(180, 220, 230)
            line.line.fill.background()

        # Step title and description
        txBox = slide.shapes.add_textbox(
            Inches(1.6), y + Inches(0.05), Inches(10.5), Inches(0.8)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        run_t = p.add_run()
        run_t.text = title
        run_t.font.size = Pt(22)
        run_t.font.bold = True
        run_t.font.color.rgb = DARK_BLUE
        run_t.font.name = FONT_NAME

        run_d = p.add_run()
        run_d.text = f"  \u2014  {desc}"
        run_d.font.size = Pt(18)
        run_d.font.color.rgb = MEDIUM_GRAY
        run_d.font.name = FONT_NAME


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5: Requirements
# ══════════════════════════════════════════════════════════════════════════
def create_slide_5():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Requirements", 5)

    requirements = [
        "Classifier must run in real-time during the demo (not pre-recorded)",
        "Data must come from the BioRadio (not simulated for final demo)",
        "You must collect your own training data (no pre-made datasets)",
        "Any ML framework and any control target are allowed",
    ]

    start_y = Inches(1.8)

    for i, req in enumerate(requirements):
        y = start_y + i * Inches(1.2)

        # Card background
        card = add_rounded_rect_card(
            slide, Inches(0.8), y, Inches(11.5), Inches(0.95),
            fill_color=WHITE, border_color=RGBColor(200, 215, 230),
            border_width=Pt(1)
        )

        # Teal left indicator
        indicator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), y,
            Inches(0.08), Inches(0.95)
        )
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = TEAL
        indicator.line.fill.background()

        # Check mark circle
        check = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.15), y + Inches(0.2),
            Inches(0.5), Inches(0.5)
        )
        check.fill.solid()
        check.fill.fore_color.rgb = TEAL
        check.line.fill.background()

        check_tf = check.text_frame
        p_check = check_tf.paragraphs[0]
        p_check.text = "\u2713"
        p_check.font.size = Pt(22)
        p_check.font.bold = True
        p_check.font.color.rgb = WHITE
        p_check.font.name = FONT_NAME
        p_check.alignment = PP_ALIGN.CENTER
        check_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Requirement text
        txBox = slide.shapes.add_textbox(
            Inches(1.9), y + Inches(0.18), Inches(10.0), Inches(0.6)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = req
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.font.name = FONT_NAME


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6: Deliverables
# ══════════════════════════════════════════════════════════════════════════
def create_slide_6():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Deliverables", 6)

    deliverables = [
        {
            "number": "1",
            "title": "Live Demo",
            "desc": "Working system that reads biosignals, runs a trained classifier, and controls something",
            "icon": "\u25B6",
        },
        {
            "number": "2",
            "title": "Presentation (~5 min)",
            "desc": "Explain your approach, results, and what you learned",
            "icon": "\u2605",
        },
    ]

    for i, d in enumerate(deliverables):
        y = Inches(1.8) + i * Inches(2.4)

        # Large card
        card = add_rounded_rect_card(
            slide, Inches(0.8), y, Inches(11.5), Inches(2.0),
            fill_color=WHITE, border_color=TEAL, border_width=Pt(2)
        )

        # Number badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.5),
            Inches(0.9), Inches(0.9)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = DARK_BLUE
        badge.line.fill.background()

        badge_tf = badge.text_frame
        p_badge = badge_tf.paragraphs[0]
        p_badge.text = d["number"]
        p_badge.font.size = Pt(32)
        p_badge.font.bold = True
        p_badge.font.color.rgb = WHITE
        p_badge.font.name = FONT_NAME
        p_badge.alignment = PP_ALIGN.CENTER
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Title
        txBox = slide.shapes.add_textbox(
            Inches(2.5), y + Inches(0.35), Inches(9.0), Inches(0.6)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = d["title"]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.font.name = FONT_NAME

        # Description
        txBox2 = slide.shapes.add_textbox(
            Inches(2.5), y + Inches(1.0), Inches(9.0), Inches(0.7)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = d["desc"]
        p2.font.size = Pt(18)
        p2.font.color.rgb = MEDIUM_GRAY
        p2.font.name = FONT_NAME


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7: Judging Rubric
# ══════════════════════════════════════════════════════════════════════════
def create_slide_7():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Judging Rubric (100 pts + 10 bonus)", 7)

    categories = [
        ("ML Classifier Performance", "30", TEAL),
        ("Creativity & Novelty", "25", RGBColor(70, 130, 180)),
        ("User Experience & System Design", "25", RGBColor(100, 149, 237)),
        ("Presentation Quality", "20", MEDIUM_BLUE),
    ]

    start_y = Inches(1.7)
    bar_max_width = Inches(7.5)

    for i, (name, pts, color) in enumerate(categories):
        y = start_y + i * Inches(1.15)

        # Category name
        txBox = slide.shapes.add_textbox(
            Inches(0.8), y, Inches(5.0), Inches(0.5)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.font.name = FONT_NAME

        # Score bar background
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), y + Inches(0.5),
            bar_max_width, Inches(0.4)
        )
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = LIGHT_GRAY
        bar_bg.line.fill.background()
        bar_bg.adjustments[0] = 0.4

        # Score bar fill
        fill_width = int(bar_max_width * int(pts) / 30)
        bar_fill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), y + Inches(0.5),
            fill_width, Inches(0.4)
        )
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = color
        bar_fill.line.fill.background()
        bar_fill.adjustments[0] = 0.4

        # Points label
        txBox2 = slide.shapes.add_textbox(
            Inches(8.6), y + Inches(0.45), Inches(2.0), Inches(0.5)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"{pts} pts"
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = color
        p2.font.name = FONT_NAME

    # Bonus section
    bonus_y = start_y + 4 * Inches(1.15) + Inches(0.3)
    bonus_card = add_rounded_rect_card(
        slide, Inches(0.8), bonus_y, Inches(11.5), Inches(0.8),
        fill_color=RGBColor(255, 250, 235),
        border_color=ACCENT_ORANGE, border_width=Pt(2)
    )

    txBox3 = slide.shapes.add_textbox(
        Inches(1.2), bonus_y + Inches(0.15), Inches(10.5), Inches(0.5)
    )
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run_star = p3.add_run()
    run_star.text = "\u2605 "
    run_star.font.size = Pt(20)
    run_star.font.color.rgb = ACCENT_ORANGE
    run_star.font.name = FONT_NAME

    run_bonus = p3.add_run()
    run_bonus.text = "Up to +10 Bonus Points "
    run_bonus.font.size = Pt(20)
    run_bonus.font.bold = True
    run_bonus.font.color.rgb = DARK_BLUE
    run_bonus.font.name = FONT_NAME

    run_detail = p3.add_run()
    run_detail.text = "(multi-modal signals, cross-user generalization, adaptive learning, technical depth)"
    run_detail.font.size = Pt(16)
    run_detail.font.color.rgb = MEDIUM_GRAY
    run_detail.font.name = FONT_NAME


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8: Scoring — ML Classifier
# ══════════════════════════════════════════════════════════════════════════
def create_slide_8():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Scoring Breakdown \u2014 ML Classifier (30 pts)", 8)

    criteria = [
        "Number of distinct classes",
        "Reported accuracy / confusion matrix",
        "Responsiveness (low latency)",
        "Handling transitions and resting states",
        "Feature engineering quality",
    ]

    # Subtitle
    txBox_sub = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(8.0), Inches(0.5)
    )
    tf_sub = txBox_sub.text_frame
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Judges look for:"
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = MEDIUM_GRAY
    p_sub.font.italic = True
    p_sub.font.name = FONT_NAME

    for i, criterion in enumerate(criteria):
        y = Inches(2.2) + i * Inches(0.95)

        # Bullet indicator
        indicator = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.9), y + Inches(0.1),
            Inches(0.35), Inches(0.35)
        )
        indicator.fill.solid()
        indicator.fill.fore_color.rgb = TEAL
        indicator.line.fill.background()
        indicator.adjustments[0] = 0.3

        ind_tf = indicator.text_frame
        p_ind = ind_tf.paragraphs[0]
        p_ind.text = "\u2713"
        p_ind.font.size = Pt(16)
        p_ind.font.bold = True
        p_ind.font.color.rgb = WHITE
        p_ind.font.name = FONT_NAME
        p_ind.alignment = PP_ALIGN.CENTER
        ind_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Text
        txBox = slide.shapes.add_textbox(
            Inches(1.5), y + Inches(0.05), Inches(10.5), Inches(0.45)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = criterion
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_TEXT
        p.font.name = FONT_NAME


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9: Scoring — Creativity & UX
# ══════════════════════════════════════════════════════════════════════════
def create_slide_9():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Scoring \u2014 Creativity & User Experience", 9)

    # Left column: Creativity
    left_card = add_rounded_rect_card(
        slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.0),
        fill_color=WHITE, border_color=TEAL, border_width=Pt(2)
    )

    # Creativity header
    header_bar_l = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.7),
        Inches(5.8), Inches(0.65)
    )
    header_bar_l.fill.solid()
    header_bar_l.fill.fore_color.rgb = DARK_BLUE
    header_bar_l.line.fill.background()

    txBox_h1 = slide.shapes.add_textbox(
        Inches(0.9), Inches(1.78), Inches(5.2), Inches(0.5)
    )
    tf_h1 = txBox_h1.text_frame
    p_h1 = tf_h1.paragraphs[0]
    p_h1.text = "Creativity & Novelty (25 pts)"
    p_h1.font.size = Pt(22)
    p_h1.font.bold = True
    p_h1.font.color.rgb = WHITE
    p_h1.font.name = FONT_NAME

    creativity_items = [
        "Original application or approach",
        "Creative use of biosignals",
        "\"Wow factor\"",
    ]

    tf_c = add_body_textbox(
        slide, Inches(0.9), Inches(2.6), Inches(5.2), Inches(3.5)
    )
    tf_c.paragraphs[0].text = ""
    for item in creativity_items:
        add_bullet(tf_c, item, size=20, color=DARK_TEXT)

    # Right column: UX
    right_card = add_rounded_rect_card(
        slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.0),
        fill_color=WHITE, border_color=RGBColor(70, 130, 180),
        border_width=Pt(2)
    )

    header_bar_r = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.7),
        Inches(5.8), Inches(0.65)
    )
    header_bar_r.fill.solid()
    header_bar_r.fill.fore_color.rgb = DARK_BLUE
    header_bar_r.line.fill.background()

    txBox_h2 = slide.shapes.add_textbox(
        Inches(7.1), Inches(1.78), Inches(5.2), Inches(0.5)
    )
    tf_h2 = txBox_h2.text_frame
    p_h2 = tf_h2.paragraphs[0]
    p_h2.text = "User Experience (25 pts)"
    p_h2.font.size = Pt(22)
    p_h2.font.bold = True
    p_h2.font.color.rgb = WHITE
    p_h2.font.name = FONT_NAME

    ux_items = [
        "Intuitive and enjoyable to use",
        "Clear feedback to the user",
        "Handles errors gracefully",
        "Easy calibration for new users",
    ]

    tf_u = add_body_textbox(
        slide, Inches(7.1), Inches(2.6), Inches(5.2), Inches(3.5)
    )
    tf_u.paragraphs[0].text = ""
    for item in ux_items:
        add_bullet(tf_u, item, size=20, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10: Schedule
# ══════════════════════════════════════════════════════════════════════════
def create_slide_10():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Schedule", 10)

    headers = ["Time Block", "Activity"]
    data = [
        ["Kickoff", "Introduction, rules, BioRadio setup, team formation"],
        ["Hours 1\u20134", "Explore signals, collect data, initial prototyping"],
        ["Hours 5\u201316", "Development (data collection, feature engineering, model training, control)"],
        ["Hours 17\u201322", "Integration, testing, refinement"],
        ["Hours 22\u201323", "Prepare presentations"],
        ["Hour 24", "Demos & presentations to judges"],
    ]

    col_widths = [Inches(2.5), Inches(8.3)]

    create_styled_table(
        slide, rows=7, cols=2,
        left=Inches(0.9), top=Inches(1.7),
        width=Inches(10.8), height=Inches(5.0),
        headers=headers, data=data, col_widths=col_widths
    )


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11: Getting Started
# ══════════════════════════════════════════════════════════════════════════
def create_slide_11():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Getting Started", 11)

    steps = [
        ("Install:", "conda env create -f environment.yml && conda activate hackathon"),
        ("Launch GUI:", "python -m src.hackathon_gui --mock"),
        ("Connect,", "configure channels, start streaming"),
        ("Stream to LSL:", 'Check "Stream to LSL" to send data to your scripts'),
        ("Build", "your classifier and control system!"),
    ]

    for i, (label, detail) in enumerate(steps):
        y = Inches(1.65) + i * Inches(1.05)

        # Step number
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.08),
            Inches(0.5), Inches(0.5)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = TEAL
        num_shape.line.fill.background()

        num_tf = num_shape.text_frame
        p_num = num_tf.paragraphs[0]
        p_num.text = str(i + 1)
        p_num.font.size = Pt(18)
        p_num.font.bold = True
        p_num.font.color.rgb = WHITE
        p_num.font.name = FONT_NAME
        p_num.alignment = PP_ALIGN.CENTER
        num_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Connecting line
        if i < 4:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.93), y + Inches(0.58),
                Inches(0.04), Inches(0.5)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(180, 220, 230)
            line.line.fill.background()

        # Text
        txBox = slide.shapes.add_textbox(
            Inches(1.5), y + Inches(0.0), Inches(11.0), Inches(0.7)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        run_label = p.add_run()
        run_label.text = f"{label} "
        run_label.font.size = Pt(20)
        run_label.font.bold = True
        run_label.font.color.rgb = DARK_BLUE
        run_label.font.name = FONT_NAME

        # For code-like steps, use monospace styling
        if i < 2:
            run_code = p.add_run()
            run_code.text = detail
            run_code.font.size = Pt(17)
            run_code.font.color.rgb = TEAL
            run_code.font.name = "Consolas"
            run_code.font.bold = False
        else:
            run_detail = p.add_run()
            run_detail.text = detail
            run_detail.font.size = Pt(19)
            run_detail.font.color.rgb = DARK_TEXT
            run_detail.font.name = FONT_NAME


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12: Tips for Success
# ══════════════════════════════════════════════════════════════════════════
def create_slide_12():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Tips for Success", 12)

    tips = [
        ("Start collecting data early",
         "your ML is only as good as your data"),
        ("Keep your classifier simple at first",
         "2-class SVM that works > 10-class DNN that doesn't"),
        ("Use mock mode (--mock) to develop",
         "without hardware"),
        ("Budget time for integration",
         "it always takes longer than expected"),
        ("Have a backup plan",
         "if ambitious approach fails, have a simpler version ready"),
        ("Test with a teammate",
         "who didn't train the model (bonus points for cross-user generalization!)"),
    ]

    for i, (bold_text, normal_text) in enumerate(tips):
        y = Inches(1.6) + i * Inches(0.9)

        # Teal diamond bullet
        diamond = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND, Inches(0.8), y + Inches(0.12),
            Inches(0.22), Inches(0.22)
        )
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = TEAL
        diamond.line.fill.background()

        # Tip text
        txBox = slide.shapes.add_textbox(
            Inches(1.3), y, Inches(11.2), Inches(0.7)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        run_bold = p.add_run()
        run_bold.text = bold_text
        run_bold.font.size = Pt(19)
        run_bold.font.bold = True
        run_bold.font.color.rgb = DARK_BLUE
        run_bold.font.name = FONT_NAME

        run_sep = p.add_run()
        run_sep.text = " \u2014 "
        run_sep.font.size = Pt(19)
        run_sep.font.color.rgb = MEDIUM_GRAY
        run_sep.font.name = FONT_NAME

        run_normal = p.add_run()
        run_normal.text = normal_text
        run_normal.font.size = Pt(19)
        run_normal.font.color.rgb = MEDIUM_GRAY
        run_normal.font.name = FONT_NAME


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13: Bonus Points
# ══════════════════════════════════════════════════════════════════════════
def create_slide_13():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    content_slide_setup(slide, "Bonus Points (up to +10)", 13)

    headers = ["Bonus Category", "Points"]
    data = [
        ["Multi-modal signals (e.g., EMG + EOG combined)", "+3"],
        ["Cross-user generalization", "+3"],
        ["Adaptive / online learning", "+2"],
        ["Exceptional technical depth", "+2"],
    ]

    col_widths = [Inches(8.0), Inches(2.8)]

    create_styled_table(
        slide, rows=5, cols=2,
        left=Inches(0.9), top=Inches(1.7),
        width=Inches(10.8), height=Inches(3.8),
        headers=headers, data=data, col_widths=col_widths
    )

    # Note card at the bottom
    note_y = Inches(5.8)
    note_card = add_rounded_rect_card(
        slide, Inches(0.9), note_y, Inches(11.4), Inches(0.8),
        fill_color=RGBColor(255, 250, 235),
        border_color=ACCENT_ORANGE, border_width=Pt(2)
    )

    txBox = slide.shapes.add_textbox(
        Inches(1.3), note_y + Inches(0.15), Inches(10.5), Inches(0.5)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run_star = p.add_run()
    run_star.text = "\u2605 "
    run_star.font.size = Pt(18)
    run_star.font.color.rgb = ACCENT_ORANGE
    run_star.font.name = FONT_NAME

    run_note = p.add_run()
    run_note.text = "Bonus points are awarded on top of the 100-point base score by the judges."
    run_note.font.size = Pt(18)
    run_note.font.color.rgb = DARK_TEXT
    run_note.font.name = FONT_NAME


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14: Closing
# ══════════════════════════════════════════════════════════════════════════
def create_slide_14():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_dark_background(slide)

    # Top accent bar
    add_accent_bar(slide, top=0, height=Inches(0.1), color=TEAL)

    # Left decorative bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.0), Inches(0.35), Inches(3.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()

    # Main message
    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(11.5), Inches(1.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Build Something Amazing"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4),
        Inches(4.3), Inches(0.04)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = TEAL
    divider.line.fill.background()

    # Sub messages
    txBox2 = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.8), Inches(10.5), Inches(0.6)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Good luck! Ask mentors for help anytime."
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_TEAL
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER

    # Info placeholders
    txBox3 = slide.shapes.add_textbox(
        Inches(2.0), Inches(4.8), Inches(9.0), Inches(1.5)
    )
    tf3 = txBox3.text_frame
    tf3.word_wrap = True

    p3 = tf3.paragraphs[0]
    p3.text = ""

    info_lines = [
        "WiFi: [network name / password]",
        "Repo: [github.com/your-org/hackathon-2026]",
        "Contact: [mentor emails / Slack channel]",
    ]

    for line in info_lines:
        p_info = tf3.add_paragraph()
        p_info.text = line
        p_info.font.size = Pt(18)
        p_info.font.color.rgb = RGBColor(140, 160, 190)
        p_info.font.name = FONT_NAME
        p_info.alignment = PP_ALIGN.CENTER
        p_info.space_after = Pt(6)

    # Bottom accent bar
    add_accent_bar(slide, top=SLIDE_HEIGHT - Inches(0.1),
                   height=Inches(0.1), color=TEAL)

    add_slide_number(slide, 14)
    slide.shapes[-1].text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 120, 150)


# ══════════════════════════════════════════════════════════════════════════
# MAIN: Generate the presentation
# ══════════════════════════════════════════════════════════════════════════
def main():
    print("Generating BioRadio Hackathon 2026 Kickoff Presentation...")

    create_slide_1()
    print("  [1/14] Title slide")

    create_slide_2()
    print("  [2/14] What You'll Build")

    create_slide_3()
    print("  [3/14] Signal Types Available")

    create_slide_4()
    print("  [4/14] The Hackathon Pipeline")

    create_slide_5()
    print("  [5/14] Requirements")

    create_slide_6()
    print("  [6/14] Deliverables")

    create_slide_7()
    print("  [7/14] Judging Rubric")

    create_slide_8()
    print("  [8/14] ML Classifier Scoring")

    create_slide_9()
    print("  [9/14] Creativity & UX Scoring")

    create_slide_10()
    print("  [10/14] Schedule")

    create_slide_11()
    print("  [11/14] Getting Started")

    create_slide_12()
    print("  [12/14] Tips for Success")

    create_slide_13()
    print("  [13/14] Bonus Points")

    create_slide_14()
    print("  [14/14] Closing")

    output_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "BioRadio_Hackathon_2026_Kickoff.pptx"
    )
    prs.save(output_path)
    print(f"\nPresentation saved to: {output_path}")
    print(f"File size: {os.path.getsize(output_path) / 1024:.1f} KB")
    print("Done!")


if __name__ == "__main__":
    main()
